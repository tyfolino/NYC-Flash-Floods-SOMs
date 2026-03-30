"""
Generate a PowerPoint presentation with all figures, supplementary figures,
and movie first-frames, each with their GRL-style captions.

Usage:
    python generate_pptx.py
Output:
    NYC_flash_flood_SOMs_figures.pptx  (in the repo root)
"""

import io
import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
BASE = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(BASE, "figures")
MOVS = os.path.join(BASE, "movies")

# Ordered list of (label, png_path, caption_path)
ENTRIES = [
    # Main figures
    (
        "Figure 1",
        os.path.join(FIGS, "fig01", "fig01_som_node_weights.png"),
        os.path.join(FIGS, "fig01", "caption.txt"),
    ),
    (
        "Figure 2",
        os.path.join(FIGS, "fig02", "fig02_evsom_key_hours.png"),
        os.path.join(FIGS, "fig02", "caption.txt"),
    ),
    (
        "Figure 3",
        os.path.join(FIGS, "fig03", "fig03_evsom_composites_key_hours.png"),
        os.path.join(FIGS, "fig03", "caption.txt"),
    ),
    (
        "Figure 4",
        os.path.join(FIGS, "fig04", "fig04_stageiv_pmm.png"),
        os.path.join(FIGS, "fig04", "caption.txt"),
    ),
    # Supplementary figures
    (
        "Figure S1",
        os.path.join(FIGS, "figS01", "figS01_ffe_distributions.png"),
        os.path.join(FIGS, "figS01", "caption.txt"),
    ),
    (
        "Figure S2",
        os.path.join(FIGS, "figS02", "figS02_sammon_map.png"),
        os.path.join(FIGS, "figS02", "caption.txt"),
    ),
    (
        "Figure S3",
        os.path.join(FIGS, "figS03", "figS03_alldays_raw_composites.png"),
        os.path.join(FIGS, "figS03", "caption.txt"),
    ),
    (
        "Figure S4",
        os.path.join(FIGS, "figS04", "figS04_ffe_som.png"),
        os.path.join(FIGS, "figS04", "caption.txt"),
    ),
    (
        "Figure S5",
        os.path.join(FIGS, "figS05", "figS05_evsom_sammon_map.png"),
        os.path.join(FIGS, "figS05", "caption.txt"),
    ),
    (
        "Figure S6",
        os.path.join(FIGS, "figS06", "figS06_evsom_to_alldays_mapping.png"),
        os.path.join(FIGS, "figS06", "caption.txt"),
    ),
    (
        "Figure S7",
        os.path.join(FIGS, "figS07", "figS07_evsom_ivt_tcwv_composites.png"),
        os.path.join(FIGS, "figS07", "caption.txt"),
    ),
    (
        "Figure S8",
        os.path.join(FIGS, "figS08", "figS08_evsom_monthly_histograms.png"),
        os.path.join(FIGS, "figS08", "caption.txt"),
    ),
    (
        "Figure S9",
        os.path.join(FIGS, "figS09", "figS09_precip_histograms.png"),
        os.path.join(FIGS, "figS09", "caption.txt"),
    ),
    (
        "Figure S10",
        os.path.join(FIGS, "figS10", "figS10_evsom_tc_tracks_and_histogram.png"),
        os.path.join(FIGS, "figS10", "caption.txt"),
    ),
    # Movies — first frame of GIF used as thumbnail
    (
        "Movie S1",
        os.path.join(MOVS, "movieS01", "movieS01_evsom_combined.gif"),
        os.path.join(MOVS, "movieS01", "caption.txt"),
    ),
]

# ---------------------------------------------------------------------------
# Slide dimensions — 16:9 widescreen
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Layout constants
MARGIN = Inches(0.25)
TITLE_H = Inches(0.40)
CAPTION_H = Inches(1.35)
IMG_H = SLIDE_H - MARGIN - TITLE_H - Inches(0.05) - CAPTION_H - MARGIN

# Colors
C_TITLE_BG = RGBColor(0x1F, 0x49, 0x7D)  # dark blue
C_TITLE_FG = RGBColor(0xFF, 0xFF, 0xFF)  # white
C_BODY_BG = RGBColor(0xF2, 0xF2, 0xF2)  # light grey
C_BODY_FG = RGBColor(0x1A, 0x1A, 0x1A)  # near-black
C_SLIDE_BG = RGBColor(0xFF, 0xFF, 0xFF)  # white


def read_caption(path: str) -> str:
    """Read caption text, stripping leading line numbers like '     1→'."""
    with open(path, encoding="utf-8") as fh:
        lines = fh.readlines()
    parts = []
    for line in lines:
        # Strip the "     N→" prefix written by the Write tool
        if "→" in line:
            line = line.split("→", 1)[1]
        # Drop system-reminder noise
        if line.strip().startswith("<system-reminder>") or line.strip().startswith(
            "</system-reminder>"
        ):
            continue
        if "<system-reminder>" in line or "</system-reminder>" in line:
            continue
        parts.append(line)
    return "".join(parts).strip()


def get_image_bytes(path: str) -> bytes:
    """Return PNG bytes for a given image path (handles GIF first frame)."""
    img = Image.open(path)
    if (
        hasattr(img, "is_animated")
        and img.is_animated
        or getattr(img, "n_frames", 1) > 1
    ):
        img.seek(0)
    img = img.convert("RGBA")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def set_slide_background(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, label: str, is_movie: bool = False):
    """Dark blue title bar at top."""
    tf = slide.shapes.add_textbox(MARGIN, MARGIN, SLIDE_W - 2 * MARGIN, TITLE_H)
    tf.fill.solid()
    tf.fill.fore_color.rgb = C_TITLE_BG
    p = tf.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    suffix = "  [animation — see .gif / .mp4]" if is_movie else ""
    run.text = label + suffix
    run.font.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = C_TITLE_FG


def add_image(slide, img_buf, label: str):
    """Center the figure image between the title bar and caption area."""
    img_top = MARGIN + TITLE_H + Inches(0.05)
    img_max_w = SLIDE_W - 2 * MARGIN

    # Measure natural image aspect ratio
    img_buf.seek(0)
    pil = Image.open(img_buf)
    nat_w, nat_h = pil.size
    aspect = nat_w / nat_h

    # Fit within (img_max_w × IMG_H) preserving aspect ratio
    if aspect > (img_max_w / IMG_H):
        w = img_max_w
        h = img_max_w / aspect
    else:
        h = IMG_H
        w = IMG_H * aspect

    left = MARGIN + (img_max_w - w) / 2
    top = img_top + (IMG_H - h) / 2

    img_buf.seek(0)
    slide.shapes.add_picture(img_buf, left, top, width=w, height=h)


def add_caption(slide, caption_text: str):
    """Light-grey caption box at the bottom of the slide."""
    cap_top = SLIDE_H - MARGIN - CAPTION_H
    tf = slide.shapes.add_textbox(MARGIN, cap_top, SLIDE_W - 2 * MARGIN, CAPTION_H)
    tf.fill.solid()
    tf.fill.fore_color.rgb = C_BODY_BG
    tf.text_frame.word_wrap = True

    p = tf.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = caption_text
    run.font.size = Pt(9)
    run.font.color.rgb = C_BODY_FG


def add_title_slide(prs: Presentation):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, C_TITLE_BG)

    # Centered title text
    tf = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), SLIDE_W - Inches(2), Inches(2)
    )
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "NYC Flash Flood SOMs"
    run.font.bold = True
    run.font.size = Pt(40)
    run.font.color.rgb = C_TITLE_FG

    tf2 = slide.shapes.add_textbox(
        Inches(1), Inches(4.5), SLIDE_W - Inches(2), Inches(1)
    )
    tf2.text_frame.word_wrap = True
    p2 = tf2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Figures, Supplementary Figures, and Movies"
    run2.font.size = Pt(22)
    run2.font.color.rgb = RGBColor(0xCC, 0xD9, 0xEA)


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs)

    for label, img_path, cap_path in ENTRIES:
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, C_SLIDE_BG)

        is_movie = label.startswith("Movie")
        add_title_bar(slide, label, is_movie=is_movie)

        caption = read_caption(cap_path)
        add_caption(slide, caption)

        if os.path.exists(img_path):
            img_buf = get_image_bytes(img_path)
            add_image(slide, img_buf, label)
            print(f"  + {label}: {os.path.basename(img_path)}")
        else:
            print(f"  ! {label}: image not found at {img_path}")

    out_path = os.path.join(BASE, "NYC_flash_flood_SOMs_figures.pptx")
    prs.save(out_path)
    print(f"\nSaved: {out_path}")


if __name__ == "__main__":
    build_presentation()
